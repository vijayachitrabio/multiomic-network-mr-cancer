#!/usr/bin/env python3
"""
Script 36: Build 15-slide professional PowerPoint presentation
Multi-Omic MR — Key Results for BMC Medicine submission
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

proj  = "."
figs  = os.path.join(proj, "results/figures")
out   = os.path.join(proj, "PRESENTATION_MultiOmicMR_2026-05-07.pptx")

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)   # title bars / headings
MID_BLUE    = RGBColor(0x2E, 0x6E, 0xA6)   # sub-headings / accents
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF5)   # background boxes
TEAL        = RGBColor(0x00, 0x83, 0x88)   # tier 2 highlights
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)   # positive / confirmed
ORANGE      = RGBColor(0xD4, 0x6B, 0x00)   # caution / placeholder
RED         = RGBColor(0xA8, 0x1C, 0x1C)   # risk / warning
GREY        = RGBColor(0x55, 0x55, 0x55)   # body text
LIGHT_GREY  = RGBColor(0xF0, 0xF0, 0xF0)  # table alt rows
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=Pt(14), bold=False, color=GREY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_bullet_box(slide, items, l, t, w, h,
                   font_size=Pt(13), title=None, title_size=Pt(15),
                   bullet="▸ ", color=GREY, title_color=DARK_BLUE,
                   line_spacing=1.15):
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size  = title_size
        run.font.bold  = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        # indent
        if isinstance(item, tuple):
            indent, text = item
        else:
            indent, text = 0, item
        p.level = indent
        run = p.add_run()
        run.text = (bullet if indent == 0 else "    • ") + text
        run.font.size = font_size
        run.font.color.rgb = color
    return txBox

def title_bar(slide, title, subtitle=None):
    """Dark blue title bar at top of slide."""
    bar = add_rect(slide, 0, 0, W, Inches(1.15), fill_rgb=DARK_BLUE)
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.07), Inches(12.7), Inches(0.65),
                 font_size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.38),
                     font_size=Pt(14), color=RGBColor(0xBB, 0xD6, 0xF0), italic=True)

def footer(slide, text="Multi-Omic MR | Modhukur et al. 2026 | DRAFT"):
    add_rect(slide, 0, Inches(7.2), W, Inches(0.3), fill_rgb=DARK_BLUE)
    add_text_box(slide, text,
                 Inches(0.2), Inches(7.21), Inches(13), Inches(0.27),
                 font_size=Pt(9), color=RGBColor(0xAA, 0xCC, 0xEE),
                 align=PP_ALIGN.CENTER)

def slide_num_box(slide, n, total=15):
    add_text_box(slide, f"{n} / {total}",
                 Inches(12.5), Inches(7.18), Inches(0.8), Inches(0.3),
                 font_size=Pt(9), color=RGBColor(0xAA, 0xCC, 0xEE),
                 align=PP_ALIGN.RIGHT)

def col_header(slide, headers, lefts, top, widths, row_h=Inches(0.38)):
    for header, l, w in zip(headers, lefts, widths):
        box = add_rect(slide, l, top, w, row_h, fill_rgb=DARK_BLUE)
        add_text_box(slide, header, l+Inches(0.06), top+Inches(0.04),
                     w-Inches(0.1), row_h-Inches(0.06),
                     font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def table_row(slide, cells, lefts, top, widths, row_h=Inches(0.33),
              bg=None, font_size=Pt(10.5), bold=False, color=GREY):
    for cell, l, w in zip(cells, lefts, widths):
        if bg:
            add_rect(slide, l, top, w, row_h, fill_rgb=bg,
                     line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
        add_text_box(slide, str(cell), l+Inches(0.05), top+Inches(0.03),
                     w-Inches(0.08), row_h-Inches(0.05),
                     font_size=font_size, bold=bold, color=color,
                     align=PP_ALIGN.CENTER)

def embed_fig(slide, fname, l, t, w, h=None):
    path = os.path.join(figs, fname)
    if not os.path.exists(path):
        print(f"  ⚠ missing: {fname}")
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# gradient background simulation
add_rect(s, 0, 0, W, H, fill_rgb=RGBColor(0xF2, 0xF7, 0xFC))
add_rect(s, 0, 0, W, Inches(3.5), fill_rgb=DARK_BLUE)
add_rect(s, 0, Inches(3.5), W, Inches(0.12), fill_rgb=TEAL)

add_text_box(s,
    "Multi-Omic Mendelian Randomization Identifies\nProtein–Metabolite Networks Causally Linked\nto Breast and Endometrial Cancer",
    Inches(0.6), Inches(0.45), Inches(12.1), Inches(2.8),
    font_size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s,
    "Vijayachitra Modhukur\n"
    "Celvia OÜ, Tartu, Estonia  |  University of Tartu, Institute of Molecular and Cell Biology",
    Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
    font_size=Pt(14), color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_text_box(s,
    "May 2026  |  Draft — BMC Medicine submission",
    Inches(1), Inches(4.6), Inches(11.3), Inches(0.45),
    font_size=Pt(13), italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# key stats boxes
box_data = [
    ("701\nProteins screened", Inches(1.2)),
    ("17\nBreast cancer hits", Inches(4.1)),
    ("PPH4 ≈ 1.0\nGCKR coloc", Inches(7.0)),
    ("p = 1.5×10⁻¹¹\nSNX15 MAGMA", Inches(9.9)),
]
for lbl, lft in box_data:
    add_rect(s, lft, Inches(5.3), Inches(2.3), Inches(1.7), fill_rgb=DARK_BLUE,
             line_rgb=TEAL, line_w=Pt(2))
    add_text_box(s, lbl, lft+Inches(0.08), Inches(5.38), Inches(2.15), Inches(1.55),
                 font_size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s); slide_num_box(s, 1)
print("✓ Slide 1: Title")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — STUDY DESIGN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Study Design", "Three-phase multi-omic MR pipeline")

# Pipeline boxes
phases = [
    ("PHASE 1\nPQTL Screen",     "FinnGen R10 Olink\nN = 619 | 701 proteins\ncis-pQTL instruments\n(GRCh38)",   DARK_BLUE),
    ("PHASE 2\nProtein→Cancer",  "Breast (N=122,977)\nEndometrial (N=12,906)\nOvarian (N=25,509)\nWald ratio / IVW", MID_BLUE),
    ("PHASE 3\nMediation MR",    "56 NMR metabolites\nUK Biobank mQTLs\nProtein→Metabolite\n→Cancer paths",  TEAL),
    ("PHASE 4\nValidation",      "Steiger filtering\nER subtype analysis\nColocalization (COLOC)\nMAGMA | Druggability", RGBColor(0x2A, 0x6B, 0x3C)),
]
arrows = ["→", "→", "→"]
for i, (title, body, col) in enumerate(phases):
    l = Inches(0.25 + i * 3.25)
    add_rect(s, l, Inches(1.4), Inches(3.0), Inches(1.0), fill_rgb=col)
    add_text_box(s, title, l+Inches(0.1), Inches(1.43), Inches(2.8), Inches(0.95),
                 font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, l, Inches(2.4), Inches(3.0), Inches(1.65), fill_rgb=LIGHT_BLUE,
             line_rgb=col, line_w=Pt(1.5))
    add_text_box(s, body, l+Inches(0.1), Inches(2.45), Inches(2.8), Inches(1.55),
                 font_size=Pt(12), color=DARK_BLUE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, "→", Inches(3.25 + i*3.25 - 0.1), Inches(1.65),
                     Inches(0.4), Inches(0.6), font_size=Pt(28), bold=True,
                     color=MID_BLUE, align=PP_ALIGN.CENTER)

# Key numbers row
add_rect(s, Inches(0.25), Inches(4.25), Inches(12.8), Inches(0.05), fill_rgb=TEAL)
add_bullet_box(s,
    ["FDR < 0.05 threshold for Phase 2 hits (Benjamini-Hochberg across 1,882 protein–cancer pairs)",
     "Colocalization PPH4 > 0.7 = strong support  |  MAGMA Bonferroni P < 2.85×10⁻⁶",
     "All instruments: cis-pQTLs within 1 Mb of gene body  |  GRCh38 coordinates  |  F-stat > 10"],
    Inches(0.4), Inches(4.45), Inches(12.5), Inches(1.35),
    font_size=Pt(12.5), color=GREY)

footer(s); slide_num_box(s, 2)
print("✓ Slide 2: Study Design")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 17 SIGNIFICANT HITS (FOREST PLOT)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Proteome-Wide MR Screen: 17 Significant Hits",
          "FDR < 0.05 across 1,882 protein–cancer pairs (701 proteins × 3 cancers)")

# Forest plot figure — left
embed_fig(s, "fig1_phase2_forest.png", Inches(0.2), Inches(1.25), Inches(7.2))

# Key stats — right
add_rect(s, Inches(7.6), Inches(1.3), Inches(5.5), Inches(5.8),
         fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_w=Pt(1))
add_text_box(s, "Key Statistics", Inches(7.75), Inches(1.38), Inches(5.2), Inches(0.45),
             font_size=Pt(16), bold=True, color=DARK_BLUE)

stats = [
    ("16", "breast cancer hits"),
    ("1", "endometrial cancer hit (ABO)"),
    ("0", "ovarian cancer hits"),
    ("13 / 17", "single-SNP Wald ratio"),
    ("4 / 17", "multi-SNP IVW (ABO, KLB, PM20D1, IL34)"),
    ("Min FDR", "8×10⁻¹³  (SNX15, protective)"),
    ("Max FDR", "0.044  (TSPAN8)"),
]
for idx, (val, lbl) in enumerate(stats):
    top = Inches(1.9 + idx * 0.64)
    bg = LIGHT_GREY if idx % 2 == 0 else WHITE
    add_rect(s, Inches(7.65), top, Inches(5.4), Inches(0.58), fill_rgb=bg)
    add_text_box(s, val, Inches(7.72), top+Inches(0.07), Inches(1.3), Inches(0.45),
                 font_size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, lbl, Inches(9.1), top+Inches(0.07), Inches(3.9), Inches(0.45),
                 font_size=Pt(12), color=GREY)

footer(s); slide_num_box(s, 3)
print("✓ Slide 3: Forest plot + stats")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TOP HITS TABLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Top 17 Hits: Summary Table",
          "Ordered by FDR; OR = odds ratio per 1 SD increase in protein; Tier based on multi-method evidence")

# Column definitions
cols  = ["Protein", "Cancer", "OR (95% CI)", "FDR", "Method", "F-stat", "Tier"]
lefts = [Inches(x) for x in [0.15, 1.65, 2.75, 4.55, 5.4, 6.55, 7.25]]
widths= [Inches(x) for x in [1.45, 1.05, 1.75, 0.8,  1.1,  0.65, 0.65]]

col_header(s, cols, lefts, Inches(1.22), widths)

# Data rows (top 17 proteins ordered by FDR)
rows = [
    ("SNX15",     "Breast", "0.917 (0.897–0.937)", "8×10⁻¹³", "Wald", "321", "T2★"),
    ("EFNA1",     "Breast", "1.136 (1.095–1.178)", "1×10⁻¹¹", "Wald", "396", "T3"),
    ("FGF5",      "Breast", "0.960 (0.947–0.973)", "1×10⁻¹¹", "Wald", "456", "T4"),
    ("SWAP70",    "Breast", "0.942 (0.918–0.966)", "2×10⁻⁴",  "Wald", "311", "T3"),
    ("TNFRSF6B",  "Breast", "1.059 (1.026–1.093)", "2×10⁻³",  "Wald", "374", "T2"),
    ("ITIH3",     "Breast", "0.965 (0.948–0.982)", "3×10⁻³",  "Wald", "289", "T3"),
    ("ABO",       "Endo",   "1.049 (1.023–1.076)", "3×10⁻³",  "IVW",  "—",   "T2"),
    ("KLB",       "Breast", "0.972 (0.957–0.987)", "3×10⁻³",  "IVW",  "—",   "T3"),
    ("PM20D1",    "Breast", "1.031 (1.010–1.053)", "5×10⁻³",  "IVW",  "—",   "T2★"),
    ("UMOD",      "Breast", "0.981 (0.969–0.993)", "7×10⁻³",  "Wald", "298", "T3"),
    ("IL34",      "Breast", "1.028 (1.008–1.049)", "0.016",   "IVW",  "—",   "T3"),
    ("APOE",      "Breast", "1.018 (1.008–1.028)", "0.022",   "Wald", "321", "T2"),
    ("INHBB",     "Breast", "0.974 (0.956–0.992)", "0.024",   "Wald", "267", "T3"),
    ("CGREF1",    "Breast", "0.972 (0.952–0.992)", "0.030",   "Wald", "345", "T4"),
    ("ATRAID",    "Breast", "0.972 (0.951–0.994)", "0.037",   "Wald", "312", "T4"),
    ("FGFR4",     "Breast", "0.988 (0.979–0.997)", "0.041",   "Wald", "410", "T4"),
    ("TSPAN8",    "Breast", "1.018 (1.001–1.036)", "0.044",   "Wald", "378", "T4"),
]
tier_cols = {"T2★": GREEN, "T2": TEAL, "T3": MID_BLUE, "T4": GREY}
for idx, row in enumerate(rows):
    top = Inches(1.6 + idx * 0.33)
    bg = LIGHT_GREY if idx % 2 == 0 else WHITE
    tier = row[6]
    tc = tier_cols.get(tier, GREY)
    # draw row cells
    for cell, l, w in zip(row[:-1], lefts[:-1], widths[:-1]):
        add_rect(s, l, top, w, Inches(0.32), fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.3))
        add_text_box(s, cell, l+Inches(0.04), top+Inches(0.03),
                     w-Inches(0.06), Inches(0.27),
                     font_size=Pt(9.5), color=DARK_BLUE if idx < 3 else GREY,
                     bold=(idx < 3), align=PP_ALIGN.CENTER)
    # tier cell
    add_rect(s, lefts[6], top, widths[6], Inches(0.32), fill_rgb=tc)
    add_text_box(s, tier, lefts[6]+Inches(0.04), top+Inches(0.03),
                 widths[6]-Inches(0.06), Inches(0.27),
                 font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# legend
add_text_box(s, "★ MR + MAGMA Bonferroni  |  T2 = Well-supported  |  T3 = Partially supported  |  T4 = Suggestive (MR only)",
             Inches(0.15), Inches(7.05), Inches(13), Inches(0.25),
             font_size=Pt(9.5), italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Extend column headers area (7 → all)
add_rect(s, lefts[6], Inches(1.22), widths[6]+Inches(5.7), Inches(0.38), fill_rgb=DARK_BLUE)
# remaining cols
extra_cols = ["MAGMA p", "ER pattern", "Drugs"]
extra_lefts= [Inches(7.95), Inches(9.05), Inches(11.45)]
extra_widths=[Inches(1.05), Inches(2.35), Inches(1.55)]
for h, l, w in zip(extra_cols, extra_lefts, extra_widths):
    add_text_box(s, h, l+Inches(0.05), Inches(1.25), w-Inches(0.08), Inches(0.33),
                 font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

extra_data = [
    ("1.5×10⁻¹¹★", "Both subtypes",   "0"),
    ("0.002",       "ER-neg specific", "0"),
    ("0.084",       "ER-neg specific", "0"),
    ("6.8×10⁻⁵",   "ER-neg specific", "0"),
    ("0.373",       "Both subtypes",   "0"),
    ("3.3×10⁻⁵",   "NS",              "0"),
    ("0.239",       "—",               "0"),
    ("0.004",       "ER-neg specific", "3"),
    ("1.4×10⁻⁶★",  "Both subtypes",   "0"),
    ("0.019",       "ER-neg specific", "0"),
    ("2.6×10⁻⁴",   "ER-neg specific", "0"),
    ("0.004",       "ER-neg specific", "0"),
    ("0.017",       "ER-neg specific", "0"),
    ("0.192",       "ER-neg specific", "0"),
    ("0.072",       "Both/NS",         "0"),
    ("0.382",       "ER-pos",          "17"),
    ("0.288",       "ER-pos",          "0"),
]
for idx, (m, er, dr) in enumerate(extra_data):
    top = Inches(1.6 + idx * 0.33)
    bg = LIGHT_GREY if idx % 2 == 0 else WHITE
    for val, l, w in zip([m, er, dr], extra_lefts, extra_widths):
        add_rect(s, l, top, w, Inches(0.32), fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.3))
        fc = GREEN if '★' in val else (ORANGE if dr != '0' and val == dr else GREY)
        add_text_box(s, val, l+Inches(0.04), top+Inches(0.03),
                     w-Inches(0.06), Inches(0.27),
                     font_size=Pt(9.5), color=fc, align=PP_ALIGN.CENTER)

footer(s); slide_num_box(s, 4)
print("✓ Slide 4: Hits table")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STEIGER + INSTRUMENT QUALITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Instrument Quality & Steiger Directionality",
          "All 17 instruments explain more variance in protein than in cancer outcome")

embed_fig(s, "sfig3_steiger_r2.png", Inches(0.2), Inches(1.25), Inches(6.8))

add_bullet_box(s,
    ["0 / 17 SNPs reversed by Steiger filtering",
     "All instruments: r²_protein >> r²_cancer (diagonal plot)",
     "Confirms protein → cancer direction, not reverse causation",
     "",
     "F-statistics (instrument strength):",
     (1, "All 17 hits: F-stat > 260 (well above F > 10 threshold)"),
     (1, "Strongest: FGF5 (F = 456), FGFR4 (F = 410), EFNA1 (F = 396)"),
     (1, "Multi-SNP proteins: IVW used; MR-Egger not feasible (n < 3 SNPs)"),
     "",
     "Sensitivity: per-SNP Wald ratios for 4 multi-SNP proteins",
     (1, "8 / 8 per-SNP estimates directionally consistent with IVW"),
     (1, "ABO, KLB, PM20D1, IL34 — all both SNPs agree"),
    ],
    Inches(7.2), Inches(1.3), Inches(5.9), Inches(5.5),
    font_size=Pt(13), color=GREY)

# result box
add_rect(s, Inches(7.2), Inches(6.3), Inches(5.9), Inches(0.75),
         fill_rgb=GREEN, line_rgb=GREEN)
add_text_box(s, "✓  No evidence of reverse causation in any of the 17 hits",
             Inches(7.35), Inches(6.38), Inches(5.6), Inches(0.6),
             font_size=Pt(13.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s); slide_num_box(s, 5)
print("✓ Slide 5: Steiger")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ER SUBTYPE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "ER Subtype Analysis: Predominantly ER-Negative Specificity",
          "Exploratory analysis — ER+ (GCST90018758) vs ER- (GCST90018759)")

embed_fig(s, "fig2_er_subtype_forest.png", Inches(0.2), Inches(1.25), Inches(7.5))

add_bullet_box(s,
    ["9 / 16 breast proteins: ER-negative specific",
     "Suggests immune/stromal rather than hormone-receptor mechanisms",
     "",
     "Notable patterns:",
     (1, "FGF5, UMOD, SWAP70, IL34, APOE, INHBB, CGREF1 → ER-neg"),
     (1, "TSPAN8 → ER-positive weighted (only exception)"),
     (1, "EFNA1 → borderline ER-neg"),
     (1, "ATRAID → direction FLIPS between subtypes (pleiotropic)"),
     (1, "SNX15, TNFRSF6B → consistent across both subtypes"),
     "",
     "Implication: therapeutic targets may be more relevant",
     "for triple-negative / ER-negative breast cancer subtype",
     "",
     "⚠ Caveat: subtype-specific multiple-testing not applied",
     "— treat as hypothesis-generating only",
    ],
    Inches(7.8), Inches(1.3), Inches(5.3), Inches(5.9),
    font_size=Pt(12.5), color=GREY)

footer(s); slide_num_box(s, 6)
print("✓ Slide 6: ER subtypes")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MAGMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "MAGMA Gene-Based Analysis: Independent GWAS Convergence",
          "SNP-wise mean model | 17,545 breast genes | 1000G EUR LD reference | MAGMA v1.10")

embed_fig(s, "fig6_magma_mr_hits.png", Inches(0.2), Inches(1.25), Inches(7.0))

# Highlight box SNX15
add_rect(s, Inches(7.4), Inches(1.3), Inches(5.7), Inches(1.55),
         fill_rgb=RGBColor(0xE8, 0xF5, 0xE9), line_rgb=GREEN, line_w=Pt(2))
add_text_box(s, "SNX15  —  Tier 2 ★",
             Inches(7.55), Inches(1.35), Inches(5.4), Inches(0.4),
             font_size=Pt(15), bold=True, color=GREEN)
add_text_box(s,
    "MR FDR = 8×10⁻¹³  |  MAGMA P = 1.47×10⁻¹¹\nRank 109 / 17,545 genes  (Bonferroni ✓)\n10 GWAS SNPs — independent of single MR instrument",
    Inches(7.55), Inches(1.75), Inches(5.4), Inches(0.95),
    font_size=Pt(12.5), color=DARK_BLUE)

# Highlight box PM20D1
add_rect(s, Inches(7.4), Inches(2.95), Inches(5.7), Inches(1.45),
         fill_rgb=RGBColor(0xE8, 0xF5, 0xE9), line_rgb=TEAL, line_w=Pt(2))
add_text_box(s, "PM20D1  —  Tier 2 ★",
             Inches(7.55), Inches(3.0), Inches(5.4), Inches(0.4),
             font_size=Pt(15), bold=True, color=TEAL)
add_text_box(s,
    "MR FDR = 0.005  |  MAGMA P = 1.43×10⁻⁶  (Bonferroni ✓)\n36 GWAS SNPs, 8 independent blocks",
    Inches(7.55), Inches(3.4), Inches(5.4), Inches(0.85),
    font_size=Pt(12.5), color=DARK_BLUE)

# Enrichment finding
add_rect(s, Inches(7.4), Inches(4.5), Inches(5.7), Inches(1.4),
         fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_w=Pt(1.5))
add_text_box(s, "MR Hits Systematically Enriched in GWAS",
             Inches(7.55), Inches(4.55), Inches(5.4), Inches(0.4),
             font_size=Pt(13.5), bold=True, color=DARK_BLUE)
add_text_box(s,
    "10 / 16 breast proteins: MAGMA p < 0.05\nBinomial test vs 5% null: P = 5.9×10⁻¹⁰\nMR hits ≠ random GWAS noise",
    Inches(7.55), Inches(4.97), Inches(5.4), Inches(0.85),
    font_size=Pt(12.5), color=DARK_BLUE)

add_text_box(s,
    "Key: MAGMA uses full GWAS (not pQTL instruments) → largely independent validation",
    Inches(7.4), Inches(6.05), Inches(5.7), Inches(0.45),
    font_size=Pt(11), italic=True, color=MID_BLUE)

footer(s); slide_num_box(s, 7)
print("✓ Slide 7: MAGMA")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — METABOLITE COLOCALIZATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Metabolite–Cancer Colocalization: GCKR Locus Dominates",
          "Exhaustive scan of 13,100 loci across 4 metabolite→breast cancer pairs (FDR < 0.05)")

embed_fig(s, "fig5_metabolite_cancer_coloc.png", Inches(0.15), Inches(1.25), Inches(7.3))

# coloc table
cheaders = ["Locus", "Lead SNP", "Metabolite", "PPH4", "Interpretation"]
clefts  = [Inches(x) for x in [7.5, 8.5, 9.65, 10.9, 11.5]]
cwidths = [Inches(x) for x in [1.0, 1.15, 1.25, 0.6,  1.65]]
col_header(s, cheaders, clefts, Inches(1.3), cwidths, row_h=Inches(0.4))

coloc_rows = [
    ("GCKR",    "rs1260326", "Glycine",    "1.000", "Tier 1 ✓"),
    ("GCKR",    "rs1260326", "Total BCAA", "1.000", "Tier 1 ✓"),
    ("GCKR",    "rs1260326", "TG/PG",      "1.000", "Tier 1 ✓"),
    ("CPS1",    "rs1047891", "Glycine",    "0.994", "Novel ✓"),
    ("CPS1",    "rs1047891", "Total BCAA", "0.989", "Novel ✓"),
    ("MLXIPL",  "rs62463430","HDL-C",      "0.967", "New finding"),
    ("MC4R",    "—",         "Total BCAA", "0.995", "Additional"),
    ("VEGFA",   "—",         "Total BCAA", "0.990", "Additional"),
]
tier1_bg = RGBColor(0xE8, 0xF5, 0xE9)
for idx, row in enumerate(coloc_rows):
    top = Inches(1.7 + idx * 0.38)
    bg = tier1_bg if row[4].startswith("Tier 1") else (LIGHT_BLUE if "Novel" in row[4] else LIGHT_GREY)
    for val, l, w in zip(row, clefts, cwidths):
        add_rect(s, l, top, w, Inches(0.36), fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.4))
        fc = GREEN if "Tier 1" in val or "Novel" in val else GREY
        add_text_box(s, val, l+Inches(0.04), top+Inches(0.03),
                     w-Inches(0.06), Inches(0.3),
                     font_size=Pt(9.5), color=fc, bold=("Tier" in val or "Novel" in val),
                     align=PP_ALIGN.CENTER)

add_text_box(s,
    "GCKR Pro446Leu (rs1260326) — known breast cancer GWAS hit\n"
    "Coloc explains causal mechanism via amino-acid / lipid metabolism\n"
    "CPS1 Thr1405Asn (rs1047891) — sub-GW threshold, novel finding",
    Inches(7.5), Inches(4.8), Inches(5.65), Inches(0.95),
    font_size=Pt(11.5), color=DARK_BLUE, italic=True)

footer(s); slide_num_box(s, 8)
print("✓ Slide 8: Colocalization")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MEDIATION MR
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Two-Step Mediation MR: Six Protein→Metabolite→Cancer Paths",
          "Product-of-coefficients; delta method SE; all paths via GCKR or CPS1 metabolite instruments")

embed_fig(s, "fig4_mediation_paths.png", Inches(0.2), Inches(1.25), Inches(7.2))

# mediation table
mheaders = ["Protein", "→ Metabolite", "→ Cancer", "Indirect β", "p-value", "Support"]
mlefts  = [Inches(x) for x in [7.45, 8.55, 9.75, 10.65, 11.55, 12.2]]
mwidths = [Inches(x) for x in [1.05, 1.15, 0.85,  0.85,  0.6,   0.95]]
col_header(s, mheaders, mlefts, Inches(1.28), mwidths, row_h=Inches(0.38))

med_rows = [
    ("APOE",    "Glycine",    "Breast", "0.0031", "0.021", "Partial★"),
    ("TNFRSF6B","Total BCAA", "Breast", "0.0028", "0.019", "Partial★"),
    ("IL34",    "Total BCAA", "Breast", "0.0024", "0.001", "MR only"),
    ("EFNA1",   "Total BCAA", "Breast", "0.0019", "0.044", "MR only"),
    ("PM20D1",  "Glycine",    "Breast", "0.0017", "0.038", "Partial"),
    ("INHBB",   "Total BCAA", "Breast", "0.0015", "0.049", "MR only"),
]
for idx, row in enumerate(med_rows):
    top = Inches(1.66 + idx * 0.42)
    bg = RGBColor(0xE8,0xF5,0xE9) if "★" in row[5] else LIGHT_GREY if idx%2==0 else WHITE
    for val, l, w in zip(row, mlefts, mwidths):
        add_rect(s, l, top, w, Inches(0.4), fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.4))
        add_text_box(s, val, l+Inches(0.04), top+Inches(0.04),
                     w-Inches(0.06), Inches(0.33),
                     font_size=Pt(10), color=GREEN if "★" in val else GREY,
                     align=PP_ALIGN.CENTER)

add_text_box(s,
    "★ Partial coloc support (PPH4 > 0.5)\n"
    "All 6 paths directionally consistent\n"
    "No path has PPH4 > 0.7 across all three causal legs\n"
    "→ Interpret as hypothesis-generating, not definitive mediation",
    Inches(7.45), Inches(4.5), Inches(5.7), Inches(1.2),
    font_size=Pt(12), italic=True, color=ORANGE)

footer(s); slide_num_box(s, 9)
print("✓ Slide 9: Mediation MR")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PATHWAY + DRUGGABILITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Pathway Enrichment & Druggability",
          "gprofiler2 on 17 proteins (background = 701) | OpenTargets Platform v4")

embed_fig(s, "sfig6_gprofiler_dotplot.png", Inches(0.15), Inches(1.25), Inches(6.9))

# Pathway summary
add_rect(s, Inches(7.2), Inches(1.3), Inches(5.85), Inches(2.4),
         fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_w=Pt(1.5))
add_text_box(s, "Pathway Enrichment (38 significant terms)",
             Inches(7.35), Inches(1.35), Inches(5.6), Inches(0.4),
             font_size=Pt(14), bold=True, color=DARK_BLUE)
add_bullet_box(s,
    ["Dominant axis: FGFR4–KLB (betaKlotho) signalling",
     "Co-enriched: IGF1R / PI3K / AKT / MAPK pathway",
     "FGF ligand binding, insulin receptor signalling",
     "Reactome: FGFR4 ligand binding, FRS-mediated cascade",
    ],
    Inches(7.35), Inches(1.75), Inches(5.6), Inches(1.8),
    font_size=Pt(12.5), color=DARK_BLUE)

# Druggability table
add_text_box(s, "OpenTargets Druggability",
             Inches(7.2), Inches(3.8), Inches(5.85), Inches(0.4),
             font_size=Pt(14), bold=True, color=DARK_BLUE)
drug_rows = [
    ("FGFR4",    "17 drugs", "Approved (SM)", "ORANTINIB Phase 3", "⚠ Paradox"),
    ("KLB",      "3 drugs",  "Advanced Clin.", "FAZPILODEMAB Ph2", "Target"),
    ("APOE",     "—",        "Structure w/lig","—",                "Novel"),
    ("Other 14", "0 drugs",  "—",             "—",                "All novel"),
]
dheaders = ["Protein", "N drugs", "Tractability", "Top drug", "Status"]
dlefts  = [Inches(x) for x in [7.25, 8.35, 9.35, 10.7, 12.3]]
dwidths = [Inches(x) for x in [1.05, 0.95, 1.3,   1.55, 0.85]]
col_header(s, dheaders, dlefts, Inches(4.25), dwidths, row_h=Inches(0.36))
for idx, row in enumerate(drug_rows):
    top = Inches(4.61 + idx * 0.36)
    bg = RGBColor(0xFF,0xF3,0xE0) if row[0]=="FGFR4" else (LIGHT_GREY if idx%2==0 else WHITE)
    for val, l, w in zip(row, dlefts, dwidths):
        add_rect(s, l, top, w, Inches(0.34), fill_rgb=bg,
                 line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.3))
        fc = ORANGE if "Paradox" in val else (GREEN if "Target" in val else GREY)
        add_text_box(s, val, l+Inches(0.04), top+Inches(0.03),
                     w-Inches(0.06), Inches(0.28),
                     font_size=Pt(9.5), color=fc, align=PP_ALIGN.CENTER)

add_text_box(s,
    "⚠ FGFR4 paradox: MR shows OR=0.988 (protective) yet FGFR4 inhibitors are cancer drugs → careful framing needed",
    Inches(7.2), Inches(6.1), Inches(5.85), Inches(0.42),
    font_size=Pt(10.5), italic=True, color=ORANGE)

footer(s); slide_num_box(s, 10)
print("✓ Slide 10: Pathway + Druggability")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — UKB-PPP REPLICATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "UKB-PPP pQTL Validation Attempt",
          "APOE Olink NGS pQTL summary stats downloaded from Synapse | N = 33,995")

# Two columns
# Left: what we found
add_rect(s, Inches(0.2), Inches(1.3), Inches(6.2), Inches(5.6),
         fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_w=Pt(1.5))
add_text_box(s, "Instrument Confirmation", Inches(0.35), Inches(1.38),
             Inches(5.9), Inches(0.42), font_size=Pt(15), bold=True, color=DARK_BLUE)

# comparison table
comp = [
    ("Metric",         "FinnGen (N=619)",  "UKB-PPP (N=33,995)"),
    ("SNP",            "chr19:44908684",   "Same position ✓"),
    ("Effect allele",  "T (ε4)",           "C/T confirmed ✓"),
    ("Beta",           "−1.170",           "−1.012"),
    ("SE",             "0.0653",           "0.0104  (6× smaller)"),
    ("LOG10P",         "57.8",             "2,069.6  (55× stronger)"),
    ("F-statistic",    "321",              "~9,520"),
]
for idx, (a, b, c) in enumerate(comp):
    top = Inches(1.85 + idx * 0.45)
    bg = DARK_BLUE if idx == 0 else (LIGHT_GREY if idx%2==0 else WHITE)
    fc_a = WHITE if idx==0 else DARK_BLUE
    fc_b = WHITE if idx==0 else GREY
    fc_c = WHITE if idx==0 else GREEN
    for val, l, w, fc in [(a, Inches(0.25), Inches(2.0), fc_a),
                           (b, Inches(2.3),  Inches(1.95), fc_b),
                           (c, Inches(4.3),  Inches(2.0), fc_c)]:
        add_rect(s, l, top, w, Inches(0.43), fill_rgb=bg,
                 line_rgb=RGBColor(0xAA,0xAA,0xAA) if idx>0 else bg, line_w=Pt(0.4))
        add_text_box(s, val, l+Inches(0.05), top+Inches(0.04),
                     w-Inches(0.08), Inches(0.35),
                     font_size=Pt(10), bold=(idx==0), color=fc, align=PP_ALIGN.CENTER)

# Right: result
add_rect(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(2.8),
         fill_rgb=RGBColor(0xFF,0xF8,0xE1), line_rgb=ORANGE, line_w=Pt(2))
add_text_box(s, "Formal Replication: Blocked", Inches(6.95), Inches(1.38),
             Inches(6.0), Inches(0.45), font_size=Pt(15), bold=True, color=ORANGE)
add_text_box(s,
    "rs429358 (APOE ε4) absent from breast GWAS\n(GCST90018757)\n\n"
    "Cause: APOE haplotype block excluded during\nGWAS quality control (common in large GWAS)\n\n"
    "→ 41 proxy variants found in overlapping positions",
    Inches(6.95), Inches(1.85), Inches(6.0), Inches(2.1),
    font_size=Pt(12.5), color=DARK_BLUE)

add_rect(s, Inches(6.8), Inches(4.2), Inches(6.3), Inches(2.5),
         fill_rgb=RGBColor(0xE8,0xF5,0xE9), line_rgb=GREEN, line_w=Pt(2))
add_text_box(s, "Proxy Analysis: Direction Consistent ✓", Inches(6.95), Inches(4.28),
             Inches(6.0), Inches(0.45), font_size=Pt(15), bold=True, color=GREEN)
add_text_box(s,
    "Best proxy: rs571415134 (chr19:45,141,737)\n"
    "F-statistic = 132  |  LOG10P = 29.9\n\n"
    "Proxy OR = 1.77 (0.93–3.37),  p = 0.085\n"
    "FinnGen OR = 1.018 (1.008–1.028)\n\n"
    "Direction: POSITIVE (both) ✓\n"
    "→ Directional triangulation only (proxy ≠ formal replication)",
    Inches(6.95), Inches(4.75), Inches(6.0), Inches(1.85),
    font_size=Pt(12.5), color=DARK_BLUE)

footer(s); slide_num_box(s, 11)
print("✓ Slide 11: UKB-PPP Replication")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OBSERVATIONAL TRIANGULATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Observational Triangulation in UKB-PPP",
          "Unadjusted case-control comparison | N = 52,995 | Olink Explore 3072 platform")

embed_fig(s, "fig3_triangulation.png", Inches(0.2), Inches(1.25), Inches(7.5))

add_bullet_box(s,
    ["10 / 17 MR-predicted directions supported observationally",
     "(59% agreement; binomial P < 0.05 vs 50% null)",
     "",
     "Strongest observational support:",
     (1, "EFNA1: cases higher +0.062 NPX, p = 1.85×10⁻¹⁷"),
     (1, "SNX15: cases lower (consistent with MR protective signal)"),
     (1, "APOE: directional agreement"),
     "",
     "⚠ Important caveats:",
     (1, "Unadjusted — no covariate correction"),
     (1, "Confounding present (BMI, age, medication)"),
     (1, "UKB-PPP includes incident + prevalent cases"),
     (1, "Treat as triangulation only, not causal validation"),
     "",
     "7 / 17 directional inconsistencies — consistent with",
     "expected noise from confounding in observational data",
    ],
    Inches(7.8), Inches(1.3), Inches(5.3), Inches(5.9),
    font_size=Pt(12.5), color=GREY)

footer(s); slide_num_box(s, 12)
print("✓ Slide 12: Observational triangulation")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — EVIDENCE TIER TABLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Evidence Tier Summary",
          "Multi-method triangulation framework — from MR-only to full colocalization support")

# Tier descriptions
tiers = [
    (GREEN,    "TIER 1",  "Robust",
     "Metabolite-cancer loci",
     "GCKR locus → Gly / BCAA / TG-by-PG → Breast  |  PPH4 ≈ 1.000\n"
     "CPS1 locus → Gly / BCAA → Breast  |  PPH4 = 0.989–0.994\n"
     "MLXIPL → HDL-C → Breast  |  PPH4 = 0.967"),

    (TEAL,     "TIER 2",  "Well-supported",
     "MR + MAGMA Bonferroni or coloc",
     "SNX15: MR FDR=8×10⁻¹³ + MAGMA p=1.5×10⁻¹¹ (Bonf ✓)\n"
     "PM20D1: MR FDR=0.005 + MAGMA p=1.4×10⁻⁶ (Bonf ✓)\n"
     "ABO→Endometrial: MR + per-SNP consistency + known biology\n"
     "APOE→Breast: MR + MAGMA p=0.004 + coloc moderate + observational\n"
     "TNFRSF6B→Breast: MR + both ER subtypes"),

    (MID_BLUE, "TIER 3",  "Partially supported",
     "MR + nominal MAGMA or observational",
     "ITIH3, SWAP70, IL34, EFNA1, KLB, UMOD, INHBB\n"
     "All have MAGMA p < 0.05 and/or observational directional agreement"),

    (GREY,     "TIER 4",  "Suggestive",
     "MR only — no MAGMA or coloc support",
     "ATRAID, FGF5, CGREF1, TSPAN8, FGFR4\n"
     "Single-SNP Wald ratio only  |  Require deCODE/UKB-PPP replication"),
]

for idx, (col, tier, label, subtitle, desc) in enumerate(tiers):
    top = Inches(1.28 + idx * 1.45)
    # colour bar
    add_rect(s, Inches(0.2), top, Inches(1.15), Inches(1.35), fill_rgb=col)
    add_text_box(s, f"{tier}\n{label}", Inches(0.22), top+Inches(0.15),
                 Inches(1.1), Inches(1.1),
                 font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # content
    add_rect(s, Inches(1.35), top, Inches(11.8), Inches(1.35),
             fill_rgb=LIGHT_GREY if idx%2==0 else WHITE,
             line_rgb=col, line_w=Pt(1.5))
    add_text_box(s, subtitle, Inches(1.5), top+Inches(0.06),
                 Inches(11.5), Inches(0.38),
                 font_size=Pt(13), bold=True, color=col)
    add_text_box(s, desc, Inches(1.5), top+Inches(0.44),
                 Inches(11.5), Inches(0.85),
                 font_size=Pt(11.5), color=GREY)

footer(s); slide_num_box(s, 13)
print("✓ Slide 13: Evidence tiers")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — LIMITATIONS & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
title_bar(s, "Limitations & Future Work",
          "Honest assessment of current evidence and priorities for strengthening the paper")

# Limitations column
add_rect(s, Inches(0.2), Inches(1.3), Inches(6.2), Inches(5.7),
         fill_rgb=RGBColor(0xFF,0xF3,0xF3), line_rgb=RED, line_w=Pt(1.5))
add_text_box(s, "⚠  Current Limitations",
             Inches(0.35), Inches(1.38), Inches(6.0), Inches(0.42),
             font_size=Pt(16), bold=True, color=RED)
add_bullet_box(s,
    ["FinnGen pQTL panel N = 619\n   → 13/17 single-SNP Wald ratios\n   → MR-Egger not feasible",
     "No full cross-platform pQTL replication\n   (UKB-PPP APOE: instrument confirmed;\n    formal replication blocked by GWAS coverage)",
     "No three-leg protein–metabolite–cancer coloc\n   (best: PPH4 = 0.51 for APOE→Gly→Breast)",
     "ER subtype analysis: no multiple-testing\n   correction → hypothesis-generating only",
     "MVMR for BMI / CRP confounding not done\n   (especially for KLB and IL34)",
     "Bidirectional MR (cancer → protein)\n   not completed",
    ],
    Inches(0.35), Inches(1.88), Inches(6.0), Inches(4.8),
    font_size=Pt(12), color=RGBColor(0x6B, 0x1C, 0x1C))

# Future work column
add_rect(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(5.7),
         fill_rgb=RGBColor(0xE8,0xF5,0xE9), line_rgb=GREEN, line_w=Pt(1.5))
add_text_box(s, "✓  Priority Next Steps",
             Inches(6.95), Inches(1.38), Inches(6.1), Inches(0.42),
             font_size=Pt(16), bold=True, color=GREEN)
add_bullet_box(s,
    ["deCODE pQTL replication (N=35,559)\n   → free at decode.com/summarydata\n   → scripts 24–26 already written",
     "Protein-side colocalization for SNX15\n   → top hit needs full FinnGen regional pQTL",
     "MVMR for BMI/CRP/lipids\n   → especially KLB (FGF21/adiposity axis)",
     "Full two-step mediation for\n   FGF23, TNFSF11/RANKL, INHBB",
     "Bidirectional MR: cancer → protein",
     "Full reference list (all [REF] tags)",
     "BMC Medicine formatting\n   → abstract ≤350 words | text ≤5,000 words",
    ],
    Inches(6.95), Inches(1.88), Inches(6.1), Inches(4.8),
    font_size=Pt(12), color=RGBColor(0x1A, 0x5C, 0x2A))

footer(s); slide_num_box(s, 14)
print("✓ Slide 14: Limitations & Future Work")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=WHITE)
add_rect(s, 0, 0, W, Inches(1.15), fill_rgb=DARK_BLUE)
add_text_box(s, "Conclusions & Key Messages",
             Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.95),
             font_size=Pt(26), bold=True, color=WHITE)

# 5 key messages
msgs = [
    (GREEN,    "1",
     "SNX15 and PM20D1: strongest evidence",
     "MR FDR < 0.01  +  MAGMA Bonferroni (p < 2.85×10⁻⁶)  +  independent GWAS methods\n→ Novel breast cancer protein candidates with convergent support from two orthogonal approaches"),
    (TEAL,     "2",
     "GCKR locus: robust metabolic mechanism",
     "PPH4 ≈ 1.000 for glycine, BCAA, TG/PG → breast cancer\n→ Amino-acid/lipid metabolism at GCKR causally linked; mechanistically explains known GWAS hit"),
    (MID_BLUE, "3",
     "ER-negative specificity: therapeutic implication",
     "9/16 breast cancer signals are ER-negative specific\n→ Targets relevant to triple-negative / harder-to-treat breast cancer subtype"),
    (ORANGE,   "4",
     "UKB-PPP confirms instrument quality for APOE",
     "LOG10P = 2,070 in N=33,995; 55× more precise than FinnGen\n→ FinnGen signal is real; formal replication blocked only by GWAS coverage gap"),
    (DARK_BLUE,"5",
     "Cautious triangulation story — not definitive",
     "No hit has full three-leg MR + coloc + replication\n→ Manuscript framed as discovery + convergent evidence; deCODE replication required before clinical claims"),
]
for idx, (col, num, heading, body) in enumerate(msgs):
    top = Inches(1.28 + idx * 1.18)
    add_rect(s, Inches(0.2), top, Inches(0.7), Inches(1.08), fill_rgb=col)
    add_text_box(s, num, Inches(0.22), top+Inches(0.2), Inches(0.65), Inches(0.65),
                 font_size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.95), top, Inches(12.15), Inches(1.08),
             fill_rgb=LIGHT_GREY if idx%2==0 else WHITE,
             line_rgb=col, line_w=Pt(1.5))
    add_text_box(s, heading, Inches(1.1), top+Inches(0.05),
                 Inches(12.0), Inches(0.38),
                 font_size=Pt(14), bold=True, color=col)
    add_text_box(s, body, Inches(1.1), top+Inches(0.44),
                 Inches(12.0), Inches(0.62),
                 font_size=Pt(11.5), color=GREY)

footer(s, "Multi-Omic MR | Modhukur et al. 2026 | DRAFT — Not for external distribution")
slide_num_box(s, 15)
print("✓ Slide 15: Conclusions")

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(out)
import os
size = os.path.getsize(out)
print(f"\n✓ Presentation saved: {out}")
print(f"  Size: {size/1024:.0f} KB  |  Slides: 15")
